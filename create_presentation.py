"""
Executive Leadership Presentation
"Replacing ObservePoint with Playwright — Analytics Validation Platform"
Prepared by: Director of Sales & Technology — Global
Audience:    CEO, CFO, Board Leadership
Output:      ObservePoint_Replacement_RFP.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Premium Executive Colour Palette ─────────────────────────────────────────
BG          = RGBColor(0x07, 0x0F, 0x1F)   # Near-black navy — authority
NAVY_CARD   = RGBColor(0x0D, 0x1D, 0x3A)   # Card background
ROYAL_CARD  = RGBColor(0x12, 0x2A, 0x52)   # Secondary card
GOLD        = RGBColor(0xC9, 0xA8, 0x4C)   # Rich muted gold — premium
GOLD_LIGHT  = RGBColor(0xE8, 0xCE, 0x82)   # Light gold for large numbers
GOLD_LINE   = RGBColor(0xA0, 0x82, 0x2D)   # Darker gold for lines
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
SILVER      = RGBColor(0xD8, 0xDC, 0xE8)   # Body text
SLATE       = RGBColor(0x8A, 0x96, 0xB0)   # Muted label text
EMERALD     = RGBColor(0x27, 0xA3, 0x60)   # Green / success
CRIMSON     = RGBColor(0xBF, 0x22, 0x22)   # Red / problem
AMBER       = RGBColor(0xD9, 0x7E, 0x14)   # Warning / moderate
DIVIDER_BG  = RGBColor(0x03, 0x08, 0x14)   # Section divider — near black
HIGHLIGHT   = RGBColor(0x1A, 0x52, 0x96)   # Blue highlight


# ── Core Helpers ─────────────────────────────────────────────────────────────

def rect(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.fill.background()
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    return s

def line(slide, l, t, w, h=0.02, color=GOLD_LINE):
    return rect(slide, l, t, w, h, color)

def txt(slide, text, l, t, w, h,
        size=16, bold=False, italic=False,
        color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def card(slide, l, t, w, h, fill=NAVY_CARD, border_color=None):
    r = rect(slide, l, t, w, h, fill)
    if border_color:
        r.line.color.rgb = border_color
        r.line.width = Pt(0.75)
    return r

def gold_line_top(slide, height=0.055):
    rect(slide, 0, 0, 13.33, height, GOLD)

def gold_line_bot(slide, y=7.38, height=0.055):
    rect(slide, 0, y, 13.33, height, GOLD)

def slide_badge(slide, label, l=0.45, t=0.18):
    rect(slide, l, t, 0.06, 0.3, GOLD)
    txt(slide, label.upper(), l + 0.14, t, 3.5, 0.3,
        size=9, color=SLATE, bold=True)

def kpi(slide, number, label, note, l, t, w=3.8, h=1.7,
        num_color=GOLD_LIGHT, card_fill=NAVY_CARD):
    card(slide, l, t, w, h, card_fill, GOLD_LINE)
    line(slide, l, t, w, 0.04, GOLD)
    txt(slide, number, l, t + 0.12, w, 0.85,
        size=52, bold=True, color=num_color, align=PP_ALIGN.CENTER)
    txt(slide, label, l, t + 0.95, w, 0.38,
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, note, l, t + 1.3, w, 0.32,
        size=9, color=SLATE, align=PP_ALIGN.CENTER, italic=True)


# ── Presentation Setup ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — EXECUTIVE COVER
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, BG)

# Full-width gold stripe
gold_line_top(s, 0.08)
gold_line_bot(s, 7.42, 0.08)

# Left vertical gold accent
rect(s, 0, 0, 0.07, 7.5, GOLD)

# Subtle background texture (layered dark rect)
rect(s, 0.07, 0, 13.26, 7.5, RGBColor(0x07, 0x0F, 0x1F))

# Confidential ribbon
rect(s, 9.8, 0.18, 3.4, 0.34, RGBColor(0x18, 0x30, 0x58))
txt(s, "STRICTLY CONFIDENTIAL  |  BOARD USE ONLY",
    9.85, 0.19, 3.3, 0.3, size=7.5, color=SLATE, bold=True)

# Pre-title label
txt(s, "STRATEGIC INVESTMENT PROPOSAL", 0.55, 1.05, 9.0, 0.38,
    size=11, color=GOLD, bold=True)

# Main title — two lines
txt(s, "Replacing ObservePoint", 0.55, 1.48, 11.5, 1.1,
    size=54, bold=True, color=WHITE)
txt(s, "with Playwright", 0.55, 2.52, 11.5, 1.0,
    size=54, bold=True, color=GOLD_LIGHT)

# Subtitle
line(s, 0.55, 3.62, 5.2, 0.025, GOLD)
txt(s, "Analytics Validation Platform  —  Airline Technology Division",
    0.55, 3.72, 10.5, 0.42, size=16, color=SILVER, italic=True)

# Three headline stats
stats = [
    ("$28,800", "Annual saving from Day 1"),
    ("43 Days",  "To contract end (March 30)"),
    ("POC Done", "Validated. Not a proposal — a plan."),
]
for i, (val, lbl) in enumerate(stats):
    x = 0.55 + i * 4.15
    rect(s, x, 4.38, 3.95, 0.06, GOLD)
    txt(s, val, x, 4.5,  3.95, 0.6,
        size=26, bold=True, color=GOLD_LIGHT, align=PP_ALIGN.CENTER)
    txt(s, lbl, x, 5.08, 3.95, 0.3,
        size=11, color=SLATE, align=PP_ALIGN.CENTER)

# Footer
txt(s, "Prepared by  Director of Global Sales & Technology",
    0.55, 6.9, 7.0, 0.38, size=10, color=SLATE)
txt(s, "February  2026", 11.5, 6.9, 1.8, 0.38,
    size=10, color=SLATE, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — EXECUTIVE SUMMARY (30-SECOND OVERVIEW)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
gold_line_top(s)
slide_badge(s, "Executive Summary")

txt(s, "The Situation in 30 Seconds", 0.45, 0.52, 10.0, 0.65,
    size=32, bold=True, color=WHITE)
line(s, 0.45, 1.18, 12.4, 0.02, GOLD_LINE)

cols = [
    ("THE PROBLEM",
     "ObservePoint contract ends\nMarch 30, 2026.\nCurrent cost: $2,400/month.\nNo renewal path.",
     CRIMSON),
    ("THE SOLUTION",
     "In-house Playwright platform.\nPOC already built and validated.\nCovers all 4 airline flows.\nZero vendor dependency.",
     HIGHLIGHT),
    ("THE RESULT",
     "$26,000–$28,800 saved annually.\nFull capability maintained.\nWe own the platform forever.\nDeployed before March 30.",
     EMERALD),
]
for i, (title, body, accent) in enumerate(cols):
    x = 0.45 + i * 4.27
    card(s, x, 1.32, 4.05, 4.85, NAVY_CARD, accent)
    rect(s, x, 1.32, 4.05, 0.07, accent)
    txt(s, title, x + 0.18, 1.45, 3.7, 0.4,
        size=12, bold=True, color=accent)
    line(s, x + 0.18, 1.88, 3.5, 0.02, accent)
    txt(s, body, x + 0.18, 2.0, 3.68, 3.8,
        size=14, color=SILVER, wrap=True)

rect(s, 0, 6.42, 13.33, 0.06, GOLD)
txt(s,
    "Leadership decision required:  Approve 6-week engineering engagement "
    "to complete platform and cut over before March 30 deadline.",
    0.45, 6.52, 12.4, 0.38,
    size=11, bold=True, color=GOLD, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE FINANCIAL CASE (CFO SLIDE)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
gold_line_top(s)
slide_badge(s, "Financial Impact")

txt(s, "The Numbers That Matter", 0.45, 0.52, 10.0, 0.65,
    size=32, bold=True, color=WHITE)
line(s, 0.45, 1.18, 12.4, 0.02, GOLD_LINE)

# Current cost box
card(s, 0.45, 1.38, 4.55, 3.8, RGBColor(0x1A, 0x08, 0x08), CRIMSON)
rect(s, 0.45, 1.38, 4.55, 0.07, CRIMSON)
txt(s, "CURRENT COST", 0.62, 1.52, 4.2, 0.32, size=10, bold=True, color=CRIMSON)
txt(s, "ObservePoint", 0.62, 1.86, 4.2, 0.45, size=16, color=SILVER)
txt(s, "$2,400", 0.62, 2.28, 4.2, 0.95,
    size=66, bold=True, color=RGBColor(0xFF, 0x5A, 0x5A), align=PP_ALIGN.CENTER)
txt(s, "per month", 0.62, 3.2, 4.2, 0.35,
    size=14, color=SILVER, align=PP_ALIGN.CENTER)
line(s, 0.75, 3.6, 4.1, 0.02, CRIMSON)
txt(s, "$28,800  per year", 0.62, 3.7, 4.2, 0.4,
    size=16, bold=True, color=RGBColor(0xFF, 0x5A, 0x5A), align=PP_ALIGN.CENTER)
txt(s, "Professional Tier licence\nExpires March 30, 2026",
    0.62, 4.16, 4.2, 0.75, size=11, color=SLATE, align=PP_ALIGN.CENTER)

# Arrow
txt(s, "->", 5.22, 2.8, 1.1, 0.8,
    size=42, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
txt(s, "REPLACE", 5.12, 3.58, 1.3, 0.28,
    size=9, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# Playwright cost box
card(s, 6.55, 1.38, 4.55, 3.8, RGBColor(0x05, 0x18, 0x0D), EMERALD)
rect(s, 6.55, 1.38, 4.55, 0.07, EMERALD)
txt(s, "PLAYWRIGHT PLATFORM", 6.72, 1.52, 4.2, 0.32, size=10, bold=True, color=EMERALD)
txt(s, "In-House Solution", 6.72, 1.86, 4.2, 0.45, size=16, color=SILVER)
txt(s, "$0", 6.72, 2.28, 4.2, 0.95,
    size=66, bold=True, color=RGBColor(0x5A, 0xFF, 0xA0), align=PP_ALIGN.CENTER)
txt(s, "licence cost", 6.72, 3.2, 4.2, 0.35,
    size=14, color=SILVER, align=PP_ALIGN.CENTER)
line(s, 6.85, 3.6, 4.1, 0.02, EMERALD)
txt(s, "$0–$2,400  per year", 6.72, 3.7, 4.2, 0.4,
    size=16, bold=True, color=RGBColor(0x5A, 0xFF, 0xA0), align=PP_ALIGN.CENTER)
txt(s, "Infrastructure only (CI/CD)\nOpen source — MIT licence",
    6.72, 4.16, 4.2, 0.75, size=11, color=SLATE, align=PP_ALIGN.CENTER)

# Annual saving mega card
card(s, 11.35, 1.38, 1.65, 3.8, NAVY_CARD, GOLD)
txt(s, "ANNUAL\nSAVING", 11.38, 1.55, 1.58, 0.7,
    size=9, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
# Vertical text simulation
for i, ch in enumerate("$26,400"):
    txt(s, ch, 11.38, 2.32 + i * 0.38, 1.58, 0.38,
        size=18, bold=True, color=GOLD_LIGHT, align=PP_ALIGN.CENTER)

# ROI summary bar
card(s, 0.45, 5.42, 12.4, 1.28, RGBColor(0x10, 0x20, 0x10), GOLD)
rect(s, 0.45, 5.42, 12.4, 0.06, GOLD)
roi_items = [
    ("PAYBACK PERIOD", "< 30 Days"),
    ("5-YEAR SAVING",  "$132,000+"),
    ("ENGINEERING COST", "~240 Hours"),
    ("RISK LEVEL",     "Low"),
]
for i, (lbl, val) in enumerate(roi_items):
    x = 0.75 + i * 3.1
    txt(s, lbl, x, 5.54, 2.85, 0.26, size=8, color=SLATE, bold=True)
    txt(s, val, x, 5.78, 2.85, 0.52, size=20, bold=True, color=GOLD_LIGHT)
    if i < 3:
        rect(s, x + 2.85, 5.55, 0.04, 0.95, GOLD_LINE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SECTION DIVIDER: THE CHALLENGE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DIVIDER_BG)

rect(s, 0, 0, 0.12, 7.5, GOLD)
txt(s, "01", 0.4, 0.9, 3.0, 2.5,
    size=160, bold=True, color=RGBColor(0x18, 0x28, 0x44))
txt(s, "THE\nCHALLENGE", 0.4, 1.8, 8.0, 2.5,
    size=62, bold=True, color=WHITE)
line(s, 0.4, 4.42, 6.5, 0.04, GOLD)
txt(s, "Why we cannot stand still — and why this is an opportunity.",
    0.4, 4.56, 9.5, 0.5, size=18, color=SILVER, italic=True)

rect(s, 0, 7.42, 13.33, 0.08, GOLD)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — WHAT OBSERVEPOINT DOES FOR US
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
gold_line_top(s)
slide_badge(s, "Current State")

txt(s, "What ObservePoint Does Today", 0.45, 0.52, 10.0, 0.65,
    size=32, bold=True, color=WHITE)
line(s, 0.45, 1.18, 12.4, 0.02, GOLD_LINE)

txt(s,
    "ObservePoint validates that analytics data is accurate across every step of our "
    "airline digital journeys — ensuring what we report to the business is trustworthy.",
    0.45, 1.28, 12.4, 0.6, size=13, color=SILVER, italic=True)

flows = [
    ("01", "Normal\nBooking",  "8 Steps\nSearch to Confirm"),
    ("02", "MMB\nRebooking",   "6 Steps\nChange & Reconfirm"),
    ("03", "Online\nCheck-in", "6 Steps\nRetrieve to Boarding Pass"),
    ("04", "Cabin\nUpgrade",   "5 Steps\nOffer to Payment"),
]
for i, (num, name, steps) in enumerate(flows):
    x = 0.45 + i * 3.22
    card(s, x, 2.05, 3.05, 2.55, NAVY_CARD, GOLD_LINE)
    rect(s, x, 2.05, 3.05, 0.06, GOLD)
    txt(s, num, x + 0.14, 2.16, 0.55, 0.45, size=18, bold=True, color=GOLD)
    txt(s, name, x + 0.14, 2.6, 2.75, 0.75,
        size=18, bold=True, color=WHITE)
    txt(s, steps, x + 0.14, 3.35, 2.75, 0.55,
        size=12, color=SLATE)

txt(s, "At every step we validate three data sources agree:",
    0.45, 4.78, 12.4, 0.35, size=14, bold=True, color=GOLD)

sources = [
    ("GA4\nNetwork Calls", "Beacons sent to Google Analytics 4\nIntercepted at HTTP layer"),
    ("digitalData Layer\n(W3C Standard)", "window.digitalData object\nAirline's master data source"),
    ("b2t Events\n(Action Types)", "Custom tagging events per step\nWith actionType property"),
]
for i, (name, desc) in enumerate(sources):
    x = 0.45 + i * 4.27
    card(s, x, 5.22, 4.05, 1.7, ROYAL_CARD)
    txt(s, name, x + 0.15, 5.3, 3.75, 0.65,
        size=13, bold=True, color=GOLD_LIGHT)
    txt(s, desc, x + 0.15, 5.95, 3.75, 0.75,
        size=11, color=SILVER)
    if i < 2:
        txt(s, "=", x + 3.9, 5.7, 0.5, 0.6,
            size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — THE URGENCY: WHY WE MUST MOVE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
gold_line_top(s)
slide_badge(s, "The Imperative")

txt(s, "Three Reasons We Cannot Renew ObservePoint",
    0.45, 0.52, 12.0, 0.65, size=32, bold=True, color=WHITE)
line(s, 0.45, 1.18, 12.4, 0.02, GOLD_LINE)

reasons = [
    ("CONTRACT\nENDS", "March 30, 2026",
     "43 days from today. ObservePoint has confirmed no contract extension "
     "is available at current pricing. We go dark on April 1 if we do nothing.",
     CRIMSON, "DEADLINE"),
    ("COST IS\nEXCESSIVE", "$28,800 / Year",
     "The Professional tier costs $2,400/month for capabilities our "
     "engineering team can build, own, and extend at near-zero marginal cost.",
     AMBER, "FINANCIAL"),
    ("PLATFORM\nIS AGING", "No SPA Support",
     "ObservePoint was designed for traditional websites. Our airline site "
     "is a Single-Page App. SPA journey validation is the #1 complaint "
     "in ObservePoint reviews — and our most critical use case.",
     HIGHLIGHT, "TECHNICAL"),
]
for i, (label, val, body, col, tag) in enumerate(reasons):
    x = 0.45 + i * 4.27
    card(s, x, 1.35, 4.05, 4.75, NAVY_CARD, col)
    rect(s, x, 1.35, 4.05, 0.07, col)
    rect(s, x + 3.3, 1.42, 0.65, 0.22, col)
    txt(s, tag, x + 3.3, 1.42, 0.65, 0.22,
        size=7, bold=True, color=BG, align=PP_ALIGN.CENTER)
    txt(s, label, x + 0.18, 1.5, 3.68, 0.72,
        size=20, bold=True, color=col)
    txt(s, val, x + 0.18, 2.22, 3.68, 0.48,
        size=18, bold=True, color=WHITE)
    line(s, x + 0.18, 2.72, 3.4, 0.02, col)
    txt(s, body, x + 0.18, 2.85, 3.68, 2.95,
        size=12, color=SILVER)

rect(s, 0, 6.3, 13.33, 0.06, GOLD)
txt(s, "The good news: our team has already solved this. "
       "A Playwright-based replacement is built, tested, and ready to deploy.",
    0.45, 6.4, 12.4, 0.45, size=12, bold=True, color=GOLD)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — SECTION DIVIDER: OUR SOLUTION
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DIVIDER_BG)

rect(s, 0, 0, 0.12, 7.5, GOLD)
txt(s, "02", 0.4, 0.9, 3.0, 2.5,
    size=160, bold=True, color=RGBColor(0x18, 0x28, 0x44))
txt(s, "OUR\nSOLUTION", 0.4, 1.8, 8.0, 2.5,
    size=62, bold=True, color=WHITE)
line(s, 0.4, 4.42, 6.5, 0.04, GOLD)
txt(s, "Built by our team. Owned by us. Better than what we are replacing.",
    0.4, 4.56, 9.5, 0.5, size=18, color=SILVER, italic=True)
rect(s, 0, 7.42, 13.33, 0.08, GOLD)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — PROOF OF CONCEPT: ALREADY DONE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
gold_line_top(s)
slide_badge(s, "Proof of Concept")

txt(s, "We Are Not Starting From Zero",
    0.45, 0.52, 12.0, 0.65, size=32, bold=True, color=WHITE)
line(s, 0.45, 1.18, 12.4, 0.02, GOLD_LINE)
txt(s, "The core platform is built. The following is validated and running today.",
    0.45, 1.25, 12.4, 0.38, size=13, color=SILVER, italic=True)

delivered = [
    ("GA4 Network Capture",         "Every analytics beacon intercepted at HTTP layer — GET, POST and sendBeacon"),
    ("digitalData Layer Read",      "window.digitalData captured at every step via in-browser JavaScript"),
    ("b2t Event Listener",          "Array intercept — zero events missed regardless of SPA timing"),
    ("Three-Way Validation Engine", "GA4 vs digitalData vs b2t — pass/fail per field per flow step"),
    ("SPA Route Detection",         "history.pushState intercept — every virtual page change detected"),
    ("Video Recording",             "Full HD session recording per run — no ObservePoint equivalent"),
    ("Playwright Trace Viewer",     "Step-by-step DOM + network replay for debugging — far superior to OP"),
    ("6-Sheet Excel Report",        "GA4 / digitalData / b2t / comparison / network / summary per run"),
    ("Live Web Dashboard",          "Auto-launching results dashboard — pass rate, trends, drill-down"),
]
for i, (title, desc) in enumerate(delivered):
    col = 0 if i < 5 else 1
    row = i if i < 5 else i - 5
    x = 0.45 if col == 0 else 6.92
    y = 1.72 + row * 1.05
    card(s, x, y, 6.25, 0.92, NAVY_CARD, GOLD_LINE)
    rect(s, x, y, 0.06, 0.92, EMERALD)
    txt(s, title, x + 0.2, y + 0.06, 5.9, 0.32, size=12, bold=True, color=GOLD_LIGHT)
    txt(s, desc,  x + 0.2, y + 0.4, 5.9, 0.42, size=10, color=SILVER)

rect(s, 0, 6.98, 13.33, 0.52, RGBColor(0x10, 0x20, 0x10))
rect(s, 0, 6.98, 13.33, 0.05, EMERALD)
txt(s, "The POC ran successfully. Adapting to the live airline site requires updating "
       "URL selectors and field-mapping config only — the validation engine is complete.",
    0.45, 7.06, 12.4, 0.38,
    size=11, bold=True, color=RGBColor(0x80, 0xFF, 0xB0))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — THE SPA ADVANTAGE (SIMPLIFIED FOR EXECS)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
gold_line_top(s)
slide_badge(s, "Technical Differentiator")

txt(s, "The SPA Problem — and Why We Solved It",
    0.45, 0.52, 12.0, 0.65, size=32, bold=True, color=WHITE)
line(s, 0.45, 1.18, 12.4, 0.02, GOLD_LINE)

# Context
card(s, 0.45, 1.3, 12.4, 0.75, ROYAL_CARD)
txt(s, "Our airline website is a Single-Page Application (SPA). "
       "Pages do not reload between steps — the browser silently updates content. "
       "This breaks how ObservePoint detects step boundaries and captures analytics data.",
    0.62, 1.38, 12.1, 0.58, size=12, color=SILVER)

# Two columns: OP vs PW
card(s, 0.45, 2.22, 5.95, 3.95, RGBColor(0x18, 0x06, 0x06), CRIMSON)
rect(s, 0.45, 2.22, 5.95, 0.07, CRIMSON)
txt(s, "ObservePoint — Breaks on SPAs", 0.62, 2.35, 5.6, 0.38,
    size=13, bold=True, color=CRIMSON)
op_issues = [
    "Detects page changes via full page reload only",
    "Uses fixed time delays between steps (fragile)",
    "Cannot detect history.pushState navigation",
    "Misses analytics events fired via sendBeacon",
    "Cannot wait for a specific data layer state",
    "Top-cited failure reason in customer reviews",
]
for i, issue in enumerate(op_issues):
    txt(s, f"  x   {issue}", 0.62, 2.85 + i * 0.47, 5.6, 0.42,
        size=11, color=RGBColor(0xFF, 0x88, 0x88))

card(s, 6.9, 2.22, 5.98, 3.95, RGBColor(0x05, 0x18, 0x0D), EMERALD)
rect(s, 6.9, 2.22, 5.98, 0.07, EMERALD)
txt(s, "Playwright — Built for SPAs", 7.07, 2.35, 5.6, 0.38,
    size=13, bold=True, color=EMERALD)
pw_solutions = [
    "History API intercept detects every SPA step",
    "Waits for exact analytics event — never a timer",
    "Detects pushState/replaceState navigation",
    "sendBeacon interceptor — zero missed calls",
    "Waits for specific b2t actionType to confirm step",
    "MutationObserver for modal-based steps",
]
for i, sol in enumerate(pw_solutions):
    txt(s, f"  v   {sol}", 7.07, 2.85 + i * 0.47, 5.6, 0.42,
        size=11, color=RGBColor(0x80, 0xFF, 0xB0))

# Bottom callout
rect(s, 0, 6.28, 13.33, 0.06, GOLD)
txt(s, "This is not a future promise — it is already coded and validated in the POC.",
    0.45, 6.38, 12.4, 0.38, size=12, bold=True, color=GOLD)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CAPABILITY COMPARISON (CLEAN VISUAL)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
gold_line_top(s)
slide_badge(s, "Capability Assessment")

txt(s, "ObservePoint vs Playwright — Direct Comparison",
    0.45, 0.52, 12.0, 0.65, size=32, bold=True, color=WHITE)
line(s, 0.45, 1.18, 12.4, 0.02, GOLD_LINE)

# Column headers
rect(s, 0.45, 1.3, 5.4, 0.38, NAVY_CARD)
rect(s, 5.95, 1.3, 3.45, 0.38, RGBColor(0x1A, 0x08, 0x08))
rect(s, 9.5,  1.3, 3.78, 0.38, RGBColor(0x05, 0x18, 0x0D))
txt(s, "Capability", 0.58, 1.3, 5.2, 0.38, size=10, bold=True, color=SLATE)
txt(s, "ObservePoint", 5.95, 1.3, 3.45, 0.38, size=10, bold=True, color=CRIMSON, align=PP_ALIGN.CENTER)
txt(s, "Playwright Platform", 9.5, 1.3, 3.78, 0.38, size=10, bold=True, color=EMERALD, align=PP_ALIGN.CENTER)

rows = [
    # (capability, OP status, PW status, OP+, PW+)
    ("Airline multi-step flow automation",    "Supported",        "Supported + superior auth",  True, True),
    ("GA4 event capture & validation",        "Supported",        "Supported + sendBeacon",     True, True),
    ("digitalData / data layer validation",   "Supported",        "Supported",                  True, True),
    ("b2t events on SPA sites",               "Unreliable",       "Fully solved",               False, True),
    ("SPA history.pushState detection",       "Not supported",    "Native capability",          False, True),
    ("Video recording of test sessions",      "Not available",    "Built-in HD recording",      False, True),
    ("Step-by-step trace for debugging",      "Not available",    "Playwright Trace Viewer",    False, True),
    ("Cross-browser (Firefox + Safari)",      "Chromium only",    "All 3 major engines",        False, True),
    ("Complex auth (OAuth / SSO / MFA)",      "Fragile",          "Fully supported",            False, True),
    ("Scheduled automated runs",              "Built-in SaaS",    "GitHub Actions cron",        True, True),
    ("Slack / Teams alerting",                "Built-in",         "Webhook — 1 day to build",  True, True),
    ("No-code GUI for analysts",              "Full GUI",         "YAML config templates",      True, False),
    ("Cloud-scale crawl 20K+ pages",          "Managed cloud",    "Needs infra (not required)", True, False),
]
RH = 0.36
for i, (cap, op, pw, op_ok, pw_ok) in enumerate(rows):
    bg_c = NAVY_CARD if i % 2 == 0 else RGBColor(0x10, 0x1C, 0x38)
    y = 1.7 + i * RH
    rect(s, 0.45, y, 5.4, RH - 0.03, bg_c)
    rect(s, 5.95, y, 3.45, RH - 0.03, bg_c)
    rect(s, 9.5,  y, 3.78, RH - 0.03, bg_c)
    txt(s, cap,  0.58, y, 5.2,  RH, size=10, color=SILVER)
    op_c = RGBColor(0xFF, 0x88, 0x88) if not op_ok else RGBColor(0x88, 0xCC, 0x88)
    pw_c = RGBColor(0x80, 0xFF, 0xB0) if pw_ok else GOLD
    pfx_op = "x  " if not op_ok else "v  "
    pfx_pw = "v  " if pw_ok else "~  "
    txt(s, pfx_op + op, 6.05, y, 3.25, RH, size=10, color=op_c)
    txt(s, pfx_pw + pw, 9.6,  y, 3.55, RH, size=10, color=pw_c)

# Legend
rect(s, 0.45, 6.38, 12.43, 0.28, RGBColor(0x12, 0x20, 0x40))
txt(s, "v = Supported or better    x = Gap or weakness    ~ = Partial / workaround available",
    0.6, 6.4, 12.1, 0.26, size=9, color=SLATE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — SECTION DIVIDER: THE INVESTMENT
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DIVIDER_BG)
rect(s, 0, 0, 0.12, 7.5, GOLD)
txt(s, "03", 0.4, 0.9, 3.0, 2.5,
    size=160, bold=True, color=RGBColor(0x18, 0x28, 0x44))
txt(s, "THE\nINVESTMENT", 0.4, 1.8, 8.0, 2.5,
    size=62, bold=True, color=WHITE)
line(s, 0.4, 4.42, 6.5, 0.04, GOLD)
txt(s, "Six weeks. One platform. Owned forever.",
    0.4, 4.56, 9.5, 0.5, size=18, color=SILVER, italic=True)
rect(s, 0, 7.42, 13.33, 0.08, GOLD)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — DELIVERY TIMELINE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
gold_line_top(s)
slide_badge(s, "Delivery Plan")

txt(s, "6 Weeks to Full Deployment — Before March 30",
    0.45, 0.52, 12.0, 0.65, size=32, bold=True, color=WHITE)
line(s, 0.45, 1.18, 12.4, 0.02, GOLD_LINE)

weeks = [
    ("WK 1\nFeb 16–22", "Foundation",
     "pytest structure\nYAML config per flow\nCookie governance\nSlack alerting",
     GOLD),
    ("WK 2\nFeb 23–Mar 1", "Booking &\nMMB Flows",
     "All 8 booking steps\nAll 6 MMB rebook steps\nb2t + digitalData live\n3-way comparison",
     EMERALD),
    ("WK 3\nMar 2–8", "Check-in &\nUpgrade Flows",
     "All 6 check-in steps\nAll 5 upgrade steps\nSPA wait strategy\nDashboard tabs",
     EMERALD),
    ("WK 4\nMar 9–15", "Tag Audit\n& Privacy",
     "Sitemap crawl\nTag presence report\nPre-consent scan\nGPC compliance",
     HIGHLIGHT),
    ("WK 5\nMar 16–22", "Reporting\n& Alerting",
     "Allure trend reports\nThreshold alerting\nCI/CD scheduled runs\nRegression detection",
     HIGHLIGHT),
    ("WK 6\nMar 23–30", "Cutover",
     "7-day parallel run\nParity validation\nObservePoint OFF\nFull handover",
     GOLD),
]
WW = 2.08
for i, (week, title, tasks, col) in enumerate(weeks):
    x = 0.45 + i * 2.14
    card(s, x, 1.35, WW, 1.08, RGBColor(0x0D, 0x1D, 0x3A), col)
    rect(s, x, 1.35, WW, 0.06, col)
    txt(s, week, x + 0.1, 1.42, WW-0.2, 0.62,
        size=9, bold=True, color=col)
    txt(s, title, x + 0.1, 2.0, WW-0.2, 0.38,
        size=10, bold=True, color=WHITE)
    # Gantt bar
    bar_col = col
    rect(s, x, 2.52, WW - 0.1, 0.28, bar_col)
    # Task list
    for j, task in enumerate(tasks.split("\n")):
        txt(s, f"  {task}", x + 0.05, 2.9 + j * 0.48, WW - 0.1, 0.44,
            size=9, color=SILVER)
    # Connector
    if i < 5:
        txt(s, ">", x + WW + 0.02, 2.58, 0.2, 0.28,
            size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# Milestone callouts
milestones = [
    (3.45,  "Wk 1-3\nCore flows LIVE\nMar 8"),
    (8.95,  "Wk 4-5\nFull parity\nMar 22"),
    (11.05, "Mar 30\nCUTOVER"),
]
for x, label in milestones:
    rect(s, x - 0.02, 5.02, 0.04, 0.75, GOLD)
    txt(s, label, x - 0.6, 5.72, 1.3, 0.65,
        size=8, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# Bottom ROI
card(s, 0.45, 6.38, 12.4, 0.75, NAVY_CARD, GOLD)
txt(s, "Critical Path", 0.62, 6.46, 2.5, 0.28, size=10, bold=True, color=GOLD)
txt(s, "All 4 airline flows live by Week 3 (March 8) — delivers 80% of ObservePoint value "
       "with 3 weeks to spare. Weeks 4–6 add tag audit, privacy scan, and the cutover buffer.",
    0.62, 6.68, 12.1, 0.38, size=11, color=SILVER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — HONEST GAPS (CREDIBILITY SLIDE)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
gold_line_top(s)
slide_badge(s, "Transparency")

txt(s, "Where We Lag — and Why It Does Not Matter",
    0.45, 0.52, 12.0, 0.65, size=32, bold=True, color=WHITE)
line(s, 0.45, 1.18, 12.4, 0.02, GOLD_LINE)
txt(s, "We present gaps transparently. Our assessment: none blocks the cutover.",
    0.45, 1.25, 12.4, 0.38, size=13, color=SILVER, italic=True)

gaps = [
    ("No-Code GUI",
     "ObservePoint has a point-and-click interface for non-developers.",
     "Our analytics validation is developer-owned. YAML config templates let\n"
     "analysts update test parameters without writing code.",
     "LOW IMPACT"),
    ("Cloud-Scale Crawling\n(20,000+ pages)",
     "ObservePoint's cloud scales to 20,000 pages unattended.",
     "We audited our actual usage: our largest crawl is under 3,000 pages.\n"
     "We are within Playwright's native capacity. Scaling is available via Crawlee.",
     "NOT APPLICABLE"),
    ("Multi-Team Governance\nDashboard",
     "ObservePoint has RBAC, SSO, and cross-team ownership tracking.",
     "We are a single-team user today. Allure Report on a shared server\n"
     "provides the historical dashboard we actually need.",
     "NOT APPLICABLE"),
]
for i, (gap, op_cap, mitigation, verdict) in enumerate(gaps):
    y = 1.72 + i * 1.7
    card(s, 0.45, y, 12.4, 1.58, NAVY_CARD, GOLD_LINE)
    rect(s, 0.45, y, 0.07, 1.58, AMBER)
    txt(s, gap,       0.62, y + 0.08, 3.5,  0.58, size=13, bold=True, color=WHITE)
    txt(s, op_cap,    0.62, y + 0.65, 3.5,  0.62, size=10, color=SLATE, italic=True)
    rect(s, 4.2, y + 0.08, 0.03, 1.42, GOLD_LINE)
    txt(s, "MITIGATION", 4.35, y + 0.08, 2.0, 0.26, size=8, bold=True, color=GOLD)
    txt(s, mitigation, 4.35, y + 0.36, 6.8, 0.88, size=10, color=SILVER)
    rect(s, 11.3, y + 0.08, 1.42, 0.36,
         EMERALD if "NOT" in verdict else AMBER)
    txt(s, verdict, 11.3, y + 0.08, 1.42, 0.36,
        size=8, bold=True, color=BG, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — THE DECISION
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
gold_line_top(s)
slide_badge(s, "Decision Required")

txt(s, "What We Are Asking Leadership to Approve",
    0.45, 0.52, 12.0, 0.65, size=32, bold=True, color=WHITE)
line(s, 0.45, 1.18, 12.4, 0.02, GOLD_LINE)

asks = [
    ("APPROVE",
     "6-Week Engineering Engagement",
     "Authorise 240 engineer-hours to complete the platform: "
     "all 4 airline flows, SPA handling, tag audit, privacy scan, "
     "Allure reporting, and CI/CD pipeline.",
     "Delivery by March 22, 2026",
     GOLD, GOLD_LIGHT),
    ("APPROVE",
     "Infrastructure Budget",
     "GitHub Actions compute and Allure Report shared hosting. "
     "Estimated cost $0–$200/month. This replaces $2,400/month "
     "ObservePoint licence from Day 1 of cutover.",
     "Budget: <$2,400/year vs $28,800/year today",
     EMERALD, RGBColor(0x80, 0xFF, 0xB0)),
    ("APPROVE",
     "7-Day Parallel Run (Mar 23–30)",
     "Run both ObservePoint and Playwright simultaneously "
     "for 7 days to validate full parity across all flows "
     "before decommissioning ObservePoint on March 30.",
     "Go/No-Go gate before contract end",
     HIGHLIGHT, RGBColor(0x88, 0xCC, 0xFF)),
]
for i, (badge, title, body, note, col, num_col) in enumerate(asks):
    y = 1.38 + i * 1.78
    card(s, 0.45, y, 12.4, 1.65, NAVY_CARD, col)
    rect(s, 0.45, y, 0.08, 1.65, col)
    rect(s, 0.55, y + 0.12, 0.85, 0.32, col)
    txt(s, badge, 0.55, y + 0.12, 0.85, 0.32,
        size=8, bold=True, color=BG, align=PP_ALIGN.CENTER)
    txt(s, str(i+1), 1.52, y + 0.05, 0.48, 0.55,
        size=32, bold=True, color=col)
    txt(s, title, 2.1, y + 0.1, 8.5, 0.42, size=16, bold=True, color=WHITE)
    txt(s, body,  2.1, y + 0.55, 9.2, 0.72, size=11, color=SILVER)
    txt(s, note,  2.1, y + 1.28, 6.0, 0.26, size=9, color=col, italic=True)

card(s, 0.45, 6.72, 12.4, 0.42, RGBColor(0x10, 0x20, 0x10), GOLD)
txt(s,
    "Expected outcome: ObservePoint OFF March 30  |  "
    "Playwright live April 1  |  $26,000+ saved from Year 1  |  "
    "Full airline analytics validation maintained",
    0.62, 6.78, 12.1, 0.32, size=10, bold=True, color=GOLD)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — CLOSING / THANK YOU
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DIVIDER_BG)
gold_line_top(s, 0.08)
gold_line_bot(s, 7.42, 0.08)
rect(s, 0, 0, 0.12, 7.5, GOLD)

txt(s, "Thank You", 0.5, 0.85, 9.0, 1.2,
    size=62, bold=True, color=WHITE)
txt(s, "Questions  &  Discussion", 0.5, 2.0, 10.0, 0.65,
    size=28, color=GOLD_LIGHT, italic=True)
line(s, 0.5, 2.72, 8.5, 0.04, GOLD)

summary = [
    ("43 Days",    "to contract end — we have a plan"),
    ("POC Proven", "running today, not a concept"),
    ("4 Flows",    "Booking  MMB  Check-in  Upgrade"),
    ("SPA Solved", "the problem ObservePoint cannot fix"),
    ("$26,400",    "saved in Year 1 alone"),
    ("6 Weeks",    "to full deployment with parallel run buffer"),
]
for i, (kpi_val, kpi_lbl) in enumerate(summary):
    col = 0 if i < 3 else 1
    row = i if i < 3 else i - 3
    x = 0.5 if col == 0 else 6.7
    y = 2.92 + row * 1.15
    card(s, x, y, 5.9, 0.98, NAVY_CARD, GOLD_LINE)
    rect(s, x, y, 5.9, 0.05, GOLD)
    txt(s, kpi_val, x + 0.18, y + 0.08, 2.1, 0.52,
        size=22, bold=True, color=GOLD_LIGHT)
    txt(s, kpi_lbl, x + 2.38, y + 0.18, 3.35, 0.45,
        size=12, color=SILVER)

txt(s, "Director of Global Sales & Technology  |  Digital Analytics Engineering  |  February 2026",
    0.5, 7.1, 12.5, 0.28, size=9, color=SLATE)


# ── Save ──────────────────────────────────────────────────────────────────────
output = "ObservePoint_Replacement_RFP.pptx"
prs.save(output)
print("Saved: " + output + "  |  " + str(len(prs.slides)) + " slides")
